"""Generador automatizado de la presentación PowerPoint (16:9) para Google Slides del TP2 SIA."""

import os
from pathlib import Path
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt

COLOR_FONDO = RGBColor(248, 249, 250)
COLOR_CARD = RGBColor(255, 255, 255)
COLOR_BORDE = RGBColor(226, 232, 240)
COLOR_TITULO = RGBColor(26, 54, 93)
COLOR_SUBTITULO = RGBColor(74, 85, 104)
COLOR_TEXTO = RGBColor(45, 55, 72)
COLOR_ACENTO = RGBColor(49, 130, 206)
COLOR_VERDE = RGBColor(56, 161, 105)
COLOR_TAG_BG = RGBColor(235, 248, 255)


def crear_slide_base(prs, titulo, subtitulo=None):
    """Crea una diapositiva en blanco con cabecera estandarizada y fondo limpio."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])

    fondo = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), Inches(13.333), Inches(7.5)
    )
    fondo.fill.solid()
    fondo.fill.fore_color.rgb = COLOR_FONDO
    fondo.line.fill.background()

    caja_header = slide.shapes.add_textbox(
        Inches(0.8), Inches(0.5), Inches(11.733), Inches(1.1)
    )
    tf = caja_header.text_frame
    tf.word_wrap = True
    tf.margin_left = tf.margin_top = tf.margin_right = tf.margin_bottom = 0

    p_tit = tf.paragraphs[0]
    p_tit.text = titulo
    p_tit.font.name = "Arial"
    p_tit.font.size = Pt(24)
    p_tit.font.bold = True
    p_tit.font.color.rgb = COLOR_TITULO

    if subtitulo:
        p_sub = tf.add_paragraph()
        p_sub.text = subtitulo
        p_sub.font.name = "Arial"
        p_sub.font.size = Pt(13)
        p_sub.font.color.rgb = COLOR_SUBTITULO
        p_sub.space_before = Pt(4)

    linea = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, Inches(0.8), Inches(1.6), Inches(11.733), Inches(0.02)
    )
    linea.fill.solid()
    linea.fill.fore_color.rgb = COLOR_BORDE
    linea.line.fill.background()

    return slide


def agregar_card(slide, left, top, width, height, color_bg=COLOR_CARD, color_borde=COLOR_BORDE):
    """Añade una tarjeta contenedora estilizada para agrupar gráficos, GIFs o texto."""
    card = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height)
    card.fill.solid()
    card.fill.fore_color.rgb = color_bg
    card.line.color.rgb = color_borde
    card.line.width = Pt(1)
    return card


import io
from PIL import Image

def agregar_imagen_segura(slide, ruta, left, top, width=None, height=None):
    """Inserta una imagen o GIF si existe en el disco, convirtiendo WebP a PNG si es necesario."""
    p = Path(ruta)
    if not p.exists():
        caja = slide.shapes.add_textbox(left, top, width or Inches(4), height or Inches(3))
        tf = caja.text_frame
        p_err = tf.paragraphs[0]
        p_err.text = f"[Archivo no encontrado:\n{p.name}]"
        p_err.font.size = Pt(11)
        p_err.font.color.rgb = RGBColor(229, 62, 62)
        return caja

    fuente_imagen = str(p)
    try:
        with Image.open(p) as img:
            if img.format == "WEBP":
                buffer = io.BytesIO()
                img.save(buffer, format="PNG")
                buffer.seek(0)
                fuente_imagen = buffer
    except Exception:
        pass

    if width and height:
        return slide.shapes.add_picture(fuente_imagen, left, top, width=width, height=height)
    elif width:
        return slide.shapes.add_picture(fuente_imagen, left, top, width=width)
    elif height:
        return slide.shapes.add_picture(fuente_imagen, left, top, height=height)
    return slide.shapes.add_picture(fuente_imagen, left, top)


def agregar_bullet_points(tf, puntos, font_size=13):
    """Agrega una lista de viñetas con formato consistente."""
    for idx, punto in enumerate(puntos):
        p = tf.add_paragraph() if idx > 0 or tf.paragraphs[0].text else tf.paragraphs[0]
        p.text = f"• {punto}"
        p.font.name = "Arial"
        p.font.size = Pt(font_size)
        p.font.color.rgb = COLOR_TEXTO
        p.space_before = Pt(8)


def crear_presentacion(ruta_salida="docs/presentacion/presentacion_tp2.pptx"):
    """Construye las 34 diapositivas estructuradas del trabajo práctico."""
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    root = Path(".")

    # -------------------------------------------------------------
    # 01. PORTADA
    # -------------------------------------------------------------
    slide1 = prs.slides.add_slide(prs.slide_layouts[6])
    fondo1 = slide1.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), Inches(13.333), Inches(7.5))
    fondo1.fill.solid()
    fondo1.fill.fore_color.rgb = COLOR_TITULO
    fondo1.line.fill.background()

    caja_tit1 = slide1.shapes.add_textbox(Inches(1.2), Inches(2.2), Inches(11), Inches(3.2))
    tf1 = caja_tit1.text_frame
    tf1.word_wrap = True
    p = tf1.paragraphs[0]
    p.text = "Aproximación de Imágenes con Algoritmos Genéticos"
    p.font.name = "Arial"
    p.font.size = Pt(36)
    p.font.bold = True
    p.font.color.rgb = RGBColor(255, 255, 255)

    p2 = tf1.add_paragraph()
    p2.text = "Sistemas de Inteligencia Artificial — Trabajo Práctico 2"
    p2.font.name = "Arial"
    p2.font.size = Pt(20)
    p2.font.color.rgb = RGBColor(190, 227, 248)
    p2.space_before = Pt(16)

    p3 = tf1.add_paragraph()
    p3.text = "Optimización evolutiva de operadores, selección, cruza, mutación y fenotipos"
    p3.font.name = "Arial"
    p3.font.size = Pt(14)
    p3.font.color.rgb = RGBColor(226, 232, 240)
    p3.space_before = Pt(20)

    # -------------------------------------------------------------
    # 02. EL PROBLEMA
    # -------------------------------------------------------------
    slide2 = crear_slide_base(prs, "El Problema: Aproximación Vectorial", "Objetivo y restricciones del trabajo")
    agregar_card(slide2, Inches(0.8), Inches(1.9), Inches(7.2), Inches(5.1))
    caja_txt2 = slide2.shapes.add_textbox(Inches(1.1), Inches(2.1), Inches(6.6), Inches(4.7))
    tf2 = caja_txt2.text_frame
    tf2.word_wrap = True
    puntos2 = [
        "Reconstruir una imagen objetivo mediante la superposición de figuras geométricas semitransparentes.",
        "Número fijo de figuras por individuo (longitud de cromosoma constante, ~100 genes).",
        "El orden de los genes en el individuo determina la capa de dibujado (z-order) sobre un lienzo de fondo fijo.",
        "Espacio de búsqueda masivo: parámetros continuos de coordenadas, colores y transparencias.",
        "Desafío computacional: la rasterización de fenotipos y cálculo de fitness es el cuello de botella."
    ]
    agregar_bullet_points(tf2, puntos2, 14)

    agregar_card(slide2, Inches(8.3), Inches(1.9), Inches(4.2), Inches(5.1))
    caja_lbl_obj = slide2.shapes.add_textbox(Inches(8.5), Inches(2.1), Inches(3.8), Inches(0.5))
    caja_lbl_obj.text_frame.paragraphs[0].text = "Imagen Objetivo de Referencia"
    caja_lbl_obj.text_frame.paragraphs[0].font.bold = True
    caja_lbl_obj.text_frame.paragraphs[0].font.size = Pt(13)
    caja_lbl_obj.text_frame.paragraphs[0].font.color.rgb = COLOR_TITULO
    agregar_imagen_segura(slide2, root / "resources/bron500.png", Inches(8.5), Inches(2.7), width=Inches(3.8))

    # -------------------------------------------------------------
    # 03. MODELO GENÉTICO
    # -------------------------------------------------------------
    slide3 = crear_slide_base(prs, "Modelo Genético: Genotipo, Fenotipo y Fitness", "Estructura de la solución candidata")
    anchos3 = [Inches(3.7), Inches(3.7), Inches(3.7)]
    titulos3 = ["Genotipo (Cromosoma)", "Fenotipo (Imagen)", "Función de Aptitud (Fitness)"]
    textos3 = [
        ["Vector de 100 figuras (genes).", "Cada gen codifica coordenadas (x, y), rotación/radios y RGBA de 8 bits.", "El locus fija el orden de dibujado."],
        ["Imagen rasterizada sobre fondo opaco fijo.", "Composición de capas translúcidas (alpha blending).", "Caché de fitness: solo se re-renderiza si el individuo mutó."],
        ["Distancia euclídea normalizada respecto a la imagen objetivo píxel a píxel.", "Escala en orden 1e-4 a 1e-3.", "A mayor aptitud, mayor fidelidad visual."]
    ]
    for i in range(3):
        x = Inches(0.8 + i * 4.0)
        agregar_card(slide3, x, Inches(1.9), anchos3[i], Inches(5.1))
        caja = slide3.shapes.add_textbox(x + Inches(0.2), Inches(2.1), anchos3[i] - Inches(0.4), Inches(4.7))
        tf = caja.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.text = titulos3[i]
        p.font.bold = True
        p.font.size = Pt(16)
        p.font.color.rgb = COLOR_TITULO
        agregar_bullet_points(tf, textos3[i], 13)

    # -------------------------------------------------------------
    # 04. HILO CONDUCTOR
    # -------------------------------------------------------------
    slide4 = crear_slide_base(prs, "Metodología: Hilo Conductor de Optimización", "Estrategia de selección secuencial hacia adelante (Forward Selection)")
    agregar_card(slide4, Inches(0.8), Inches(1.9), Inches(11.733), Inches(5.1))
    caja_hc = slide4.shapes.add_textbox(Inches(1.1), Inches(2.1), Inches(11.1), Inches(4.7))
    tf_hc = caja_hc.text_frame
    tf_hc.word_wrap = True
    p_hc_intro = tf_hc.paragraphs[0]
    p_hc_intro.text = "Para evitar una explosión combinatoria inabordable, se optimiza un componente a la vez en escala de 100 genes:"
    p_hc_intro.font.size = Pt(14)
    p_hc_intro.font.bold = True
    p_hc_intro.font.color.rgb = COLOR_TITULO

    etapas = [
        "Etapa 1: Inicialización — Comparación de generación aleatoria vs sesgos de color.",
        "Etapa 2: Selección de Padres — Evaluación de presión selectiva y retención de diversidad.",
        "Etapa 3: Operadores de Cruza — Análisis de disrupción de esquemas y recombinación.",
        "Etapa 4: Operadores de Mutación — Calibración de tasa y alcance de perturbación.",
        "Etapa 5: Supervivencia y Brecha — Aditiva vs Exclusiva y efecto del parámetro K.",
        "Etapa 6: Geometría de Genes — Desempeño y costo computacional según tipo de figura."
    ]
    agregar_bullet_points(tf_hc, etapas, 13)

    # -------------------------------------------------------------
    # 05. ETAPA 1: INICIALIZACIÓN - CONCEPTO
    # -------------------------------------------------------------
    slide5 = crear_slide_base(prs, "Etapa 1: Inicialización de la Población", "Estrategias de generación inicial y sesgo de color")
    agregar_card(slide5, Inches(0.8), Inches(1.9), Inches(5.7), Inches(5.1))
    caja5_1 = slide5.shapes.add_textbox(Inches(1.0), Inches(2.1), Inches(5.3), Inches(4.7))
    tf5_1 = caja5_1.text_frame
    tf5_1.word_wrap = True
    tf5_1.paragraphs[0].text = "Estrategias Evaluadas"
    tf5_1.paragraphs[0].font.bold = True
    tf5_1.paragraphs[0].font.size = Pt(16)
    tf5_1.paragraphs[0].font.color.rgb = COLOR_TITULO
    pts5_1 = [
        "Aleatoria Uniforme: Colores y coordenadas muestreados al azar en todo el lienzo.",
        "Sesgo Bounding Box: Color promedio de la caja envolvente de la figura en el objetivo.",
        "Sesgo Exact Match: Promedio exacto rasterizando la máscara del polígono u óvalo."
    ]
    agregar_bullet_points(tf5_1, pts5_1, 13)

    agregar_card(slide5, Inches(6.8), Inches(1.9), Inches(5.7), Inches(5.1))
    caja5_2 = slide5.shapes.add_textbox(Inches(7.0), Inches(2.1), Inches(5.3), Inches(4.7))
    tf5_2 = caja5_2.text_frame
    tf5_2.word_wrap = True
    tf5_2.paragraphs[0].text = "Tradeoffs Hipotéticos"
    tf5_2.paragraphs[0].font.bold = True
    tf5_2.paragraphs[0].font.size = Pt(16)
    tf5_2.paragraphs[0].font.color.rgb = COLOR_TITULO
    pts5_2 = [
        "Aleatorio: Máxima diversidad inicial, pero fitness de partida ínfimo (~1e-4) desperdiciando generaciones para fijar el fondo.",
        "Bounding Box: O(1) geométrico, salto inmediato en fitness con costo de cómputo despreciable.",
        "Exact Match: Mayor fidelidad cromática pero requiere rasterizar cada figura en gen 0."
    ]
    agregar_bullet_points(tf5_2, pts5_2, 13)

    # -------------------------------------------------------------
    # 06. ETAPA 1: GRÁFICOS
    # -------------------------------------------------------------
    slide6 = crear_slide_base(prs, "Etapa 1: Rendimiento y Convergencia", "Comparativa cuantitativa a 30 generaciones")
    agregar_card(slide6, Inches(0.8), Inches(1.9), Inches(5.7), Inches(5.1))
    caja_g6_1 = slide6.shapes.add_textbox(Inches(1.0), Inches(2.0), Inches(5.3), Inches(0.4))
    caja_g6_1.text_frame.paragraphs[0].text = "Evolución del Fitness Máximo"
    caja_g6_1.text_frame.paragraphs[0].font.bold = True
    caja_g6_1.text_frame.paragraphs[0].font.size = Pt(13)
    agregar_imagen_segura(slide6, root / "results/analysis/full_01_inicializacion/fitness_maximo.png", Inches(1.0), Inches(2.5), width=Inches(5.3))

    agregar_card(slide6, Inches(6.8), Inches(1.9), Inches(5.7), Inches(5.1))
    caja_g6_2 = slide6.shapes.add_textbox(Inches(7.0), Inches(2.0), Inches(5.3), Inches(0.4))
    caja_g6_2.text_frame.paragraphs[0].text = "Tiempo por Generación"
    caja_g6_2.text_frame.paragraphs[0].font.bold = True
    caja_g6_2.text_frame.paragraphs[0].font.size = Pt(13)
    agregar_imagen_segura(slide6, root / "results/analysis/full_01_inicializacion/tiempo_generacion.png", Inches(7.0), Inches(2.5), width=Inches(5.3))

    # -------------------------------------------------------------
    # 07. ETAPA 1: EVOLUCIÓN VISUAL (GIFs)
    # -------------------------------------------------------------
    slide7 = crear_slide_base(prs, "Etapa 1: Reconstrucción Visual", "Animación evolutiva generada por cada variante")
    var_init = [
        ("Aleatoria Uniforme", "results/analysis/full_01_inicializacion/aleatoria/img/bron500.gif"),
        ("Sesgo Bounding Box", "results/analysis/full_01_inicializacion/sesgo_bbox/img/bron500.gif"),
        ("Sesgo Exact Match", "results/analysis/full_01_inicializacion/sesgo_exact/img/bron500.gif")
    ]
    for i, (nombre, ruta) in enumerate(var_init):
        x = Inches(0.8 + i * 4.0)
        agregar_card(slide7, x, Inches(1.9), Inches(3.7), Inches(5.1))
        caja_lbl = slide7.shapes.add_textbox(x + Inches(0.1), Inches(2.1), Inches(3.5), Inches(0.4))
        caja_lbl.text_frame.paragraphs[0].text = nombre
        caja_lbl.text_frame.paragraphs[0].font.bold = True
        caja_lbl.text_frame.paragraphs[0].font.size = Pt(13)
        caja_lbl.text_frame.paragraphs[0].font.color.rgb = COLOR_TITULO
        agregar_imagen_segura(slide7, root / ruta, x + Inches(0.2), Inches(2.7), width=Inches(3.3))

    # -------------------------------------------------------------
    # 08. ETAPA 1: CONCLUSIÓN
    # -------------------------------------------------------------
    slide8 = crear_slide_base(prs, "Etapa 1: Conclusión y Selección", "Elección del método óptimo de inicialización")
    agregar_card(slide8, Inches(0.8), Inches(1.9), Inches(11.733), Inches(5.1))
    caja8 = slide8.shapes.add_textbox(Inches(1.1), Inches(2.1), Inches(11.1), Inches(4.7))
    tf8 = caja8.text_frame
    tf8.word_wrap = True
    tf8.paragraphs[0].text = "Ganador de la Etapa: Sesgo de Color (Bounding Box / Exact Match)"
    tf8.paragraphs[0].font.bold = True
    tf8.paragraphs[0].font.size = Pt(18)
    tf8.paragraphs[0].font.color.rgb = COLOR_VERDE

    pts8 = [
        "El sesgo de color logra un fitness final de 0.00065 vs 0.00013 de la inicialización aleatoria (mejora de 5x en 30 generaciones).",
        "Bounding Box alcanza prácticamente la misma convergencia que Exact Match con menor costo en inicialización.",
        "Se adopta Sesgo Bounding Box para todas las etapas siguientes, permitiendo que las figuras arranquen con tonos cromáticos pertinentes."
    ]
    agregar_bullet_points(tf8, pts8, 14)

    # -------------------------------------------------------------
    # 09. ETAPA 2: SELECCIÓN - BENCHMARK AISLADO
    # -------------------------------------------------------------
    slide9 = crear_slide_base(prs, "Etapa 2: Métodos de Selección (Evaluación Aislada)", "Análisis de la presión selectiva sobre población fija")
    agregar_card(slide9, Inches(0.8), Inches(1.9), Inches(5.7), Inches(5.1))
    caja9_txt = slide9.shapes.add_textbox(Inches(1.0), Inches(2.1), Inches(5.3), Inches(4.7))
    tf9_txt = caja9_txt.text_frame
    tf9_txt.word_wrap = True
    tf9_txt.paragraphs[0].text = "Comportamiento de Selección"
    tf9_txt.paragraphs[0].font.bold = True
    tf9_txt.paragraphs[0].font.size = Pt(16)
    tf9_txt.paragraphs[0].font.color.rgb = COLOR_TITULO
    pts9 = [
        "Ruleta / Universal: Debido a diferencias mínimas de aptitud (1e-4), la ruleta asigna probabilidades casi planas (presión selectiva nula).",
        "Torneo Determinístico: Ejerce fuerte presión seleccionando a los más aptos de cada grupo sin depender de la escala numérica.",
        "Ranking: Provee presión lineal uniforme desacoplada de la magnitud del fitness.",
        "Elite: Concentra el 100% de la masa en los primeros lugares."
    ]
    agregar_bullet_points(tf9_txt, pts9, 13)

    agregar_card(slide9, Inches(6.8), Inches(1.9), Inches(5.7), Inches(5.1))
    caja9_lbl = slide9.shapes.add_textbox(Inches(7.0), Inches(2.0), Inches(5.3), Inches(0.4))
    caja9_lbl.text_frame.paragraphs[0].text = "Frecuencia de Selección vs Ranking"
    caja9_lbl.text_frame.paragraphs[0].font.bold = True
    caja9_lbl.text_frame.paragraphs[0].font.size = Pt(13)
    agregar_imagen_segura(slide9, root / "results/analysis/op_01_seleccion/frecuencia_seleccion.png", Inches(7.0), Inches(2.5), width=Inches(5.3))

    # -------------------------------------------------------------
    # 10. ETAPA 2: SELECCIÓN - EVOLUCIÓN
    # -------------------------------------------------------------
    slide10 = crear_slide_base(prs, "Etapa 2: Selección en Ciclo Evolutivo", "Curvas generacionales de fitness y diversidad")
    agregar_card(slide10, Inches(0.8), Inches(1.9), Inches(5.7), Inches(5.1))
    caja10_1 = slide10.shapes.add_textbox(Inches(1.0), Inches(2.0), Inches(5.3), Inches(0.4))
    caja10_1.text_frame.paragraphs[0].text = "Fitness Máximo vs Generación"
    caja10_1.text_frame.paragraphs[0].font.bold = True
    caja10_1.text_frame.paragraphs[0].font.size = Pt(13)
    agregar_imagen_segura(slide10, root / "results/analysis/full_02_seleccion/fitness_maximo.png", Inches(1.0), Inches(2.5), width=Inches(5.3))

    agregar_card(slide10, Inches(6.8), Inches(1.9), Inches(5.7), Inches(5.1))
    caja10_2 = slide10.shapes.add_textbox(Inches(7.0), Inches(2.0), Inches(5.3), Inches(0.4))
    caja10_2.text_frame.paragraphs[0].text = "Diversidad Genética Normalizada"
    caja10_2.text_frame.paragraphs[0].font.bold = True
    caja10_2.text_frame.paragraphs[0].font.size = Pt(13)
    agregar_imagen_segura(slide10, root / "results/analysis/full_02_seleccion/diversidad.png", Inches(7.0), Inches(2.5), width=Inches(5.3))

    # -------------------------------------------------------------
    # 11. ETAPA 2: VISUAL (GIFs Parte 1)
    # -------------------------------------------------------------
    slide11 = crear_slide_base(prs, "Etapa 2: Reconstrucción Visual (Grupo A)", "Elite, Torneo Determinístico y Ranking")
    var_sel_a = [
        ("Elite", "results/analysis/full_02_seleccion/elite/img/bron500.gif"),
        ("Torneo Det (k=3)", "results/analysis/full_02_seleccion/torneo_det_3/img/bron500.gif"),
        ("Ranking", "results/analysis/full_02_seleccion/ranking/img/bron500.gif")
    ]
    for i, (nombre, ruta) in enumerate(var_sel_a):
        x = Inches(0.8 + i * 4.0)
        agregar_card(slide11, x, Inches(1.9), Inches(3.7), Inches(5.1))
        caja_lbl = slide11.shapes.add_textbox(x + Inches(0.1), Inches(2.1), Inches(3.5), Inches(0.4))
        caja_lbl.text_frame.paragraphs[0].text = nombre
        caja_lbl.text_frame.paragraphs[0].font.bold = True
        caja_lbl.text_frame.paragraphs[0].font.size = Pt(13)
        agregar_imagen_segura(slide11, root / ruta, x + Inches(0.2), Inches(2.7), width=Inches(3.3))

    # -------------------------------------------------------------
    # 12. ETAPA 2: VISUAL (GIFs Parte 2)
    # -------------------------------------------------------------
    slide12 = crear_slide_base(prs, "Etapa 2: Reconstrucción Visual (Grupo B)", "Ruleta, Torneo Probabilístico y Boltzmann")
    var_sel_b = [
        ("Ruleta", "results/analysis/full_02_seleccion/ruleta/img/bron500.gif"),
        ("Torneo Prob (p=0.75)", "results/analysis/full_02_seleccion/torneo_prob_75/img/bron500.gif"),
        ("Boltzmann", "results/analysis/full_02_seleccion/boltzmann/img/bron500.gif")
    ]
    for i, (nombre, ruta) in enumerate(var_sel_b):
        x = Inches(0.8 + i * 4.0)
        agregar_card(slide12, x, Inches(1.9), Inches(3.7), Inches(5.1))
        caja_lbl = slide12.shapes.add_textbox(x + Inches(0.1), Inches(2.1), Inches(3.5), Inches(0.4))
        caja_lbl.text_frame.paragraphs[0].text = nombre
        caja_lbl.text_frame.paragraphs[0].font.bold = True
        caja_lbl.text_frame.paragraphs[0].font.size = Pt(13)
        agregar_imagen_segura(slide12, root / ruta, x + Inches(0.2), Inches(2.7), width=Inches(3.3))

    # -------------------------------------------------------------
    # 13. ETAPA 2: CONCLUSIÓN
    # -------------------------------------------------------------
    slide13 = crear_slide_base(prs, "Etapa 2: Conclusión y Selección", "Elección del método óptimo de selección")
    agregar_card(slide13, Inches(0.8), Inches(1.9), Inches(11.733), Inches(5.1))
    caja13 = slide13.shapes.add_textbox(Inches(1.1), Inches(2.1), Inches(11.1), Inches(4.7))
    tf13 = caja13.text_frame
    tf13.word_wrap = True
    tf13.paragraphs[0].text = "Ganador de la Etapa: Torneo Determinístico (k=3) / Ranking"
    tf13.paragraphs[0].font.bold = True
    tf13.paragraphs[0].font.size = Pt(18)
    tf13.paragraphs[0].font.color.rgb = COLOR_VERDE

    pts13 = [
        "Torneo Determinístico ofrece una excelente presión selectiva ajustable mediante el parámetro k.",
        "Ruleta y Universal quedan descartadas para este dominio por la inviabilidad de trabajar con fitnesses que difieren en quintas decimales.",
        "Elite colapsa la diversidad prematuramente hacia copias del mismo individuo.",
        "Se traslada Torneo Determinístico a las siguientes etapas."
    ]
    agregar_bullet_points(tf13, pts13, 14)

    # -------------------------------------------------------------
    # 14. ETAPA 3: CRUZA - CONCEPTO
    # -------------------------------------------------------------
    slide14 = crear_slide_base(prs, "Etapa 3: Operadores de Cruza", "Recombinación genética y preservación de capas")
    agregar_card(slide14, Inches(0.8), Inches(1.9), Inches(11.733), Inches(5.1))
    caja14 = slide14.shapes.add_textbox(Inches(1.1), Inches(2.1), Inches(11.1), Inches(4.7))
    tf14 = caja14.text_frame
    tf14.word_wrap = True
    tf14.paragraphs[0].text = "Mecanismos y Tradeoffs"
    tf14.paragraphs[0].font.bold = True
    tf14.paragraphs[0].font.size = Pt(16)
    tf14.paragraphs[0].font.color.rgb = COLOR_TITULO
    pts14 = [
        "Un Punto: Intercambia desde una posición P hasta el final. Preserva dos grandes bloques.",
        "Dos Puntos: Intercambia el bloque intermedio [P1, P2). Preserva prefijos y sufijos de figuras contiguas.",
        "Anular: Trata al cromosoma como un anillo, permitiendo intercambiar bloques circulares que cruzan los extremos.",
        "Uniforme: Sortea locus por locus con probabilidad P. Altamente disruptivo para el orden de dibujado (z-order)."
    ]
    agregar_bullet_points(tf14, pts14, 14)

    # -------------------------------------------------------------
    # 15. ETAPA 3: CRUZA - BENCHMARK AISLADO
    # -------------------------------------------------------------
    slide15 = crear_slide_base(prs, "Etapa 3: Benchmark Aislado de Cruza", "Cantidad de genes intercambiados y disrupción por locus")
    agregar_card(slide15, Inches(0.8), Inches(1.9), Inches(5.7), Inches(5.1))
    caja15_1 = slide15.shapes.add_textbox(Inches(1.0), Inches(2.0), Inches(5.3), Inches(0.4))
    caja15_1.text_frame.paragraphs[0].text = "Genes Intercambiados (Distribución)"
    caja15_1.text_frame.paragraphs[0].font.bold = True
    caja15_1.text_frame.paragraphs[0].font.size = Pt(13)
    agregar_imagen_segura(slide15, root / "results/analysis/op_02_cruza/genes_intercambiados.png", Inches(1.0), Inches(2.5), width=Inches(5.3))

    agregar_card(slide15, Inches(6.8), Inches(1.9), Inches(5.7), Inches(5.1))
    caja15_2 = slide15.shapes.add_textbox(Inches(7.0), Inches(2.0), Inches(5.3), Inches(0.4))
    caja15_2.text_frame.paragraphs[0].text = "Probabilidad de Intercambio por Locus"
    caja15_2.text_frame.paragraphs[0].font.bold = True
    caja15_2.text_frame.paragraphs[0].font.size = Pt(13)
    agregar_imagen_segura(slide15, root / "results/analysis/op_02_cruza/frecuencia_locus.png", Inches(7.0), Inches(2.5), width=Inches(5.3))

    # -------------------------------------------------------------
    # 16. ETAPA 3: CRUZA - EVOLUCIÓN
    # -------------------------------------------------------------
    slide16 = crear_slide_base(prs, "Etapa 3: Cruza en Ciclo Evolutivo", "Convergencia y diversidad generacional")
    agregar_card(slide16, Inches(0.8), Inches(1.9), Inches(5.7), Inches(5.1))
    caja16_1 = slide16.shapes.add_textbox(Inches(1.0), Inches(2.0), Inches(5.3), Inches(0.4))
    caja16_1.text_frame.paragraphs[0].text = "Fitness Máximo vs Generación"
    caja16_1.text_frame.paragraphs[0].font.bold = True
    caja16_1.text_frame.paragraphs[0].font.size = Pt(13)
    agregar_imagen_segura(slide16, root / "results/analysis/full_03_cruza/fitness_maximo.png", Inches(1.0), Inches(2.5), width=Inches(5.3))

    agregar_card(slide16, Inches(6.8), Inches(1.9), Inches(5.7), Inches(5.1))
    caja16_2 = slide16.shapes.add_textbox(Inches(7.0), Inches(2.0), Inches(5.3), Inches(0.4))
    caja16_2.text_frame.paragraphs[0].text = "Diversidad Genética"
    caja16_2.text_frame.paragraphs[0].font.bold = True
    caja16_2.text_frame.paragraphs[0].font.size = Pt(13)
    agregar_imagen_segura(slide16, root / "results/analysis/full_03_cruza/diversidad.png", Inches(7.0), Inches(2.5), width=Inches(5.3))

    # -------------------------------------------------------------
    # 17. ETAPA 3: VISUAL (GIFs)
    # -------------------------------------------------------------
    slide17 = crear_slide_base(prs, "Etapa 3: Reconstrucción Visual", "Comparación de los 4 métodos de cruza")
    var_cruza = [
        ("Un Punto", "results/analysis/full_03_cruza/un_punto/img/bron500.gif"),
        ("Dos Puntos", "results/analysis/full_03_cruza/dos_puntos/img/bron500.gif"),
        ("Anular", "results/analysis/full_03_cruza/anular/img/bron500.gif"),
        ("Uniforme", "results/analysis/full_03_cruza/uniforme/img/bron500.gif")
    ]
    for i, (nombre, ruta) in enumerate(var_cruza):
        x = Inches(0.8 + i * 2.95)
        agregar_card(slide17, x, Inches(1.9), Inches(2.8), Inches(5.1))
        caja_lbl = slide17.shapes.add_textbox(x + Inches(0.1), Inches(2.1), Inches(2.6), Inches(0.4))
        caja_lbl.text_frame.paragraphs[0].text = nombre
        caja_lbl.text_frame.paragraphs[0].font.bold = True
        caja_lbl.text_frame.paragraphs[0].font.size = Pt(12)
        agregar_imagen_segura(slide17, root / ruta, x + Inches(0.15), Inches(2.7), width=Inches(2.5))

    # -------------------------------------------------------------
    # 18. ETAPA 3: CONCLUSIÓN
    # -------------------------------------------------------------
    slide18 = crear_slide_base(prs, "Etapa 3: Conclusión y Selección", "Elección del método óptimo de cruza")
    agregar_card(slide18, Inches(0.8), Inches(1.9), Inches(11.733), Inches(5.1))
    caja18 = slide18.shapes.add_textbox(Inches(1.1), Inches(2.1), Inches(11.1), Inches(4.7))
    tf18 = caja18.text_frame
    tf18.word_wrap = True
    tf18.paragraphs[0].text = "Ganador de la Etapa: Dos Puntos / Anular"
    tf18.paragraphs[0].font.bold = True
    tf18.paragraphs[0].font.size = Pt(18)
    tf18.paragraphs[0].font.color.rgb = COLOR_VERDE

    pts18 = [
        "Dos Puntos y Anular logran la mayor tasa de mejora sostenida al preservar bloques de capas continuas.",
        "Cruce Uniforme resulta excesivamente disruptivo para las relaciones espaciales entre figuras superpuestas.",
        "Se fija Cruce en Dos Puntos para las etapas posteriores."
    ]
    agregar_bullet_points(tf18, pts18, 14)

    # -------------------------------------------------------------
    # 19. ETAPA 4: MUTACIÓN - CONCEPTO
    # -------------------------------------------------------------
    slide19 = crear_slide_base(prs, "Etapa 4: Operadores de Mutación", "Ajuste fino local vs disrupción fenotípica")
    agregar_card(slide19, Inches(0.8), Inches(1.9), Inches(11.733), Inches(5.1))
    caja19 = slide19.shapes.add_textbox(Inches(1.1), Inches(2.1), Inches(11.1), Inches(4.7))
    tf19 = caja19.text_frame
    tf19.word_wrap = True
    tf19.paragraphs[0].text = "Métodos de Mutación Analizados"
    tf19.paragraphs[0].font.bold = True
    tf19.paragraphs[0].font.size = Pt(16)
    tf19.paragraphs[0].font.color.rgb = COLOR_TITULO
    pts19 = [
        "Gen: Muta exactamente 1 figura con probabilidad Pm. Perturbación muy conservadora.",
        "MultiGen: Muta entre 1 y M figuras sorteando cada una con probabilidad Pm. Permite ajustes locales moderados.",
        "Uniforme: Cada gen tiene probabilidad independiente Pm de mutar (media = gene_count * Pm).",
        "No Uniforme: Sorteo único: o muta el 100% de los genes o ninguno (todo o nada)."
    ]
    agregar_bullet_points(tf19, pts19, 14)

    # -------------------------------------------------------------
    # 20. ETAPA 4: MUTACIÓN - BENCHMARK AISLADO
    # -------------------------------------------------------------
    slide20 = crear_slide_base(prs, "Etapa 4: Benchmark Aislado de Mutación", "Distribución de cantidad de genes mutados por individuo")
    agregar_card(slide20, Inches(0.8), Inches(1.9), Inches(5.7), Inches(5.1))
    caja20_txt = slide20.shapes.add_textbox(Inches(1.0), Inches(2.1), Inches(5.3), Inches(4.7))
    tf20_txt = caja20_txt.text_frame
    tf20_txt.word_wrap = True
    tf20_txt.paragraphs[0].text = "Análisis de Dispersión"
    tf20_txt.paragraphs[0].font.bold = True
    tf20_txt.paragraphs[0].font.size = Pt(16)
    pts20 = [
        "Gen muta en promedio 0.5 genes.",
        "MultiGen (k=5) muta 1.5 genes en promedio, limitando la disrupción máxima a 5.",
        "MultiGen (k=15) muta ~4 genes.",
        "No Uniforme presenta alta varianza (desvío 30): salta entre 0 y 100 genes mutados en un solo ciclo."
    ]
    agregar_bullet_points(tf20_txt, pts20, 13)

    agregar_card(slide20, Inches(6.8), Inches(1.9), Inches(5.7), Inches(5.1))
    caja20_lbl = slide20.shapes.add_textbox(Inches(7.0), Inches(2.0), Inches(5.3), Inches(0.4))
    caja20_lbl.text_frame.paragraphs[0].text = "Genes Mutados (Media y Desvío Acotado)"
    caja20_lbl.text_frame.paragraphs[0].font.bold = True
    caja20_lbl.text_frame.paragraphs[0].font.size = Pt(13)
    agregar_imagen_segura(slide20, root / "results/analysis/op_03_mutacion/genes_mutados.png", Inches(7.0), Inches(2.5), width=Inches(5.3))

    # -------------------------------------------------------------
    # 21. ETAPA 4: MUTACIÓN - EVOLUCIÓN
    # -------------------------------------------------------------
    slide21 = crear_slide_base(prs, "Etapa 4: Mutación en Ciclo Evolutivo", "Convergencia y diversidad generacional")
    agregar_card(slide21, Inches(0.8), Inches(1.9), Inches(5.7), Inches(5.1))
    caja21_1 = slide21.shapes.add_textbox(Inches(1.0), Inches(2.0), Inches(5.3), Inches(0.4))
    caja21_1.text_frame.paragraphs[0].text = "Fitness Máximo vs Generación"
    caja21_1.text_frame.paragraphs[0].font.bold = True
    caja21_1.text_frame.paragraphs[0].font.size = Pt(13)
    agregar_imagen_segura(slide21, root / "results/analysis/full_04_mutacion/fitness_maximo.png", Inches(1.0), Inches(2.5), width=Inches(5.3))

    agregar_card(slide21, Inches(6.8), Inches(1.9), Inches(5.7), Inches(5.1))
    caja21_2 = slide21.shapes.add_textbox(Inches(7.0), Inches(2.0), Inches(5.3), Inches(0.4))
    caja21_2.text_frame.paragraphs[0].text = "Diversidad Genética"
    caja21_2.text_frame.paragraphs[0].font.bold = True
    caja21_2.text_frame.paragraphs[0].font.size = Pt(13)
    agregar_imagen_segura(slide21, root / "results/analysis/full_04_mutacion/diversidad.png", Inches(7.0), Inches(2.5), width=Inches(5.3))

    # -------------------------------------------------------------
    # 22. ETAPA 4: VISUAL (GIFs)
    # -------------------------------------------------------------
    slide22 = crear_slide_base(prs, "Etapa 4: Reconstrucción Visual", "Comparación de los métodos de mutación")
    var_mut = [
        ("Gen", "results/analysis/full_04_mutacion/gen/img/bron500.gif"),
        ("MultiGen (k=5)", "results/analysis/full_04_mutacion/multigen_5/img/bron500.gif"),
        ("Uniforme", "results/analysis/full_04_mutacion/uniforme/img/bron500.gif"),
        ("No Uniforme", "results/analysis/full_04_mutacion/no_uniforme/img/bron500.gif")
    ]
    for i, (nombre, ruta) in enumerate(var_mut):
        x = Inches(0.8 + i * 2.95)
        agregar_card(slide22, x, Inches(1.9), Inches(2.8), Inches(5.1))
        caja_lbl = slide22.shapes.add_textbox(x + Inches(0.1), Inches(2.1), Inches(2.6), Inches(0.4))
        caja_lbl.text_frame.paragraphs[0].text = nombre
        caja_lbl.text_frame.paragraphs[0].font.bold = True
        caja_lbl.text_frame.paragraphs[0].font.size = Pt(12)
        agregar_imagen_segura(slide22, root / ruta, x + Inches(0.15), Inches(2.7), width=Inches(2.5))

    # -------------------------------------------------------------
    # 23. ETAPA 4: CONCLUSIÓN
    # -------------------------------------------------------------
    slide23 = crear_slide_base(prs, "Etapa 4: Conclusión y Selección", "Elección del método óptimo de mutación")
    agregar_card(slide23, Inches(0.8), Inches(1.9), Inches(11.733), Inches(5.1))
    caja23 = slide23.shapes.add_textbox(Inches(1.1), Inches(2.1), Inches(11.1), Inches(4.7))
    tf23 = caja23.text_frame
    tf23.word_wrap = True
    tf23.paragraphs[0].text = "Ganador de la Etapa: MultiGen (k = 5)"
    tf23.paragraphs[0].font.bold = True
    tf23.paragraphs[0].font.size = Pt(18)
    tf23.paragraphs[0].font.color.rgb = COLOR_VERDE

    pts23 = [
        "MultiGen permite un ajuste gradual y localizado, manteniendo intacto el 90-95% del fenotipo ya adaptado.",
        "Mutación Gen (1 solo gen) progresa con excesiva lentitud.",
        "Mutación No Uniforme destruye las estructuras consolidadas cuando el evento ocurre.",
        "Se adopta MultiGen con cota máxima de 5 a 15 genes."
    ]
    agregar_bullet_points(tf23, pts23, 14)

    # -------------------------------------------------------------
    # 24. ETAPA 5: SUPERVIVENCIA - CONCEPTO
    # -------------------------------------------------------------
    slide24 = crear_slide_base(prs, "Etapa 5: Estrategias de Supervivencia", "Reemplazo generacional y brecha K")
    agregar_card(slide24, Inches(0.8), Inches(1.9), Inches(11.733), Inches(5.1))
    caja24 = slide24.shapes.add_textbox(Inches(1.1), Inches(2.1), Inches(11.1), Inches(4.7))
    tf24 = caja24.text_frame
    tf24.word_wrap = True
    tf24.paragraphs[0].text = "Dinámica de Reemplazo"
    tf24.paragraphs[0].font.bold = True
    tf24.paragraphs[0].font.size = Pt(16)
    pts24 = [
        "Supervivencia Aditiva: Padres e hijos compiten juntos en un pozo común de tamaño N + K.",
        "Supervivencia Exclusiva: La nueva generación se conforma estrictamente con hijos (con K >= N ningún padre sobrevive).",
        "Brecha generacional K: Proporción de la población renovada en cada ciclo respecto al total N.",
        "Elitismo y Caché: En aditiva, los padres superiores sobreviven sin necesidad de re-evaluar fitness."
    ]
    agregar_bullet_points(tf24, pts24, 14)

    # -------------------------------------------------------------
    # 25. ETAPA 5: BENCHMARK AISLADO DE SUPERVIVENCIA
    # -------------------------------------------------------------
    slide25 = crear_slide_base(prs, "Etapa 5: Benchmark Aislado de Supervivencia", "Preservación de padres según estrategia y brecha")
    agregar_card(slide25, Inches(0.8), Inches(1.9), Inches(5.7), Inches(5.1))
    caja25_txt = slide25.shapes.add_textbox(Inches(1.0), Inches(2.1), Inches(5.3), Inches(4.7))
    tf25_txt = caja25_txt.text_frame
    tf25_txt.word_wrap = True
    tf25_txt.paragraphs[0].text = "Efecto de K y Estrategia"
    tf25_txt.paragraphs[0].font.bold = True
    tf25_txt.paragraphs[0].font.size = Pt(16)
    pts25 = [
        "Con K = 100 (K = N), Exclusiva reemplaza el 100% de la población (0 padres sobreviven).",
        "Con K = 20, Exclusiva preserva exactamente N - K = 80 padres por diseño.",
        "En Aditiva, la cantidad de padres depende puramente de la calidad relativa de los hijos producidos."
    ]
    agregar_bullet_points(tf25_txt, pts25, 13)

    agregar_card(slide25, Inches(6.8), Inches(1.9), Inches(5.7), Inches(5.1))
    caja25_lbl = slide25.shapes.add_textbox(Inches(7.0), Inches(2.0), Inches(5.3), Inches(0.4))
    caja25_lbl.text_frame.paragraphs[0].text = "Padres Sobrevivientes en la Nueva Generación"
    caja25_lbl.text_frame.paragraphs[0].font.bold = True
    caja25_lbl.text_frame.paragraphs[0].font.size = Pt(13)
    agregar_imagen_segura(slide25, root / "results/analysis/op_04_supervivencia/supervivencia_padres.png", Inches(7.0), Inches(2.5), width=Inches(5.3))

    # -------------------------------------------------------------
    # 26. ETAPA 5: SUPERVIVENCIA - EVOLUCIÓN
    # -------------------------------------------------------------
    slide26 = crear_slide_base(prs, "Etapa 5: Supervivencia en Ciclo Evolutivo", "Monotonía del fitness y preservación")
    agregar_card(slide26, Inches(0.8), Inches(1.9), Inches(5.7), Inches(5.1))
    caja26_1 = slide26.shapes.add_textbox(Inches(1.0), Inches(2.0), Inches(5.3), Inches(0.4))
    caja26_1.text_frame.paragraphs[0].text = "Fitness Máximo vs Generación"
    caja26_1.text_frame.paragraphs[0].font.bold = True
    caja26_1.text_frame.paragraphs[0].font.size = Pt(13)
    agregar_imagen_segura(slide26, root / "results/analysis/full_05_supervivencia/fitness_maximo.png", Inches(1.0), Inches(2.5), width=Inches(5.3))

    agregar_card(slide26, Inches(6.8), Inches(1.9), Inches(5.7), Inches(5.1))
    caja26_2 = slide26.shapes.add_textbox(Inches(7.0), Inches(2.0), Inches(5.3), Inches(0.4))
    caja26_2.text_frame.paragraphs[0].text = "Diversidad Genética"
    caja26_2.text_frame.paragraphs[0].font.bold = True
    caja26_2.text_frame.paragraphs[0].font.size = Pt(13)
    agregar_imagen_segura(slide26, root / "results/analysis/full_05_supervivencia/diversidad.png", Inches(7.0), Inches(2.5), width=Inches(5.3))

    # -------------------------------------------------------------
    # 27. ETAPA 5: VISUAL (GIFs)
    # -------------------------------------------------------------
    slide27 = crear_slide_base(prs, "Etapa 5: Reconstrucción Visual", "Supervivencia Aditiva vs Exclusiva con distintos K")
    var_sup = [
        ("Aditiva (K=10)", "results/analysis/full_05_supervivencia/aditiva_k10/img/bron500.gif"),
        ("Exclusiva (K=10)", "results/analysis/full_05_supervivencia/exclusiva_k10/img/bron500.gif"),
        ("Aditiva (K=30)", "results/analysis/full_05_supervivencia/aditiva_k30/img/bron500.gif"),
        ("Exclusiva (K=30)", "results/analysis/full_05_supervivencia/exclusiva_k30/img/bron500.gif")
    ]
    for i, (nombre, ruta) in enumerate(var_sup):
        x = Inches(0.8 + i * 2.95)
        agregar_card(slide27, x, Inches(1.9), Inches(2.8), Inches(5.1))
        caja_lbl = slide27.shapes.add_textbox(x + Inches(0.1), Inches(2.1), Inches(2.6), Inches(0.4))
        caja_lbl.text_frame.paragraphs[0].text = nombre
        caja_lbl.text_frame.paragraphs[0].font.bold = True
        caja_lbl.text_frame.paragraphs[0].font.size = Pt(12)
        agregar_imagen_segura(slide27, root / ruta, x + Inches(0.15), Inches(2.7), width=Inches(2.5))

    # -------------------------------------------------------------
    # 28. ETAPA 5: CONCLUSIÓN
    # -------------------------------------------------------------
    slide28 = crear_slide_base(prs, "Etapa 5: Conclusión y Selección", "Elección de la estrategia óptima de supervivencia")
    agregar_card(slide28, Inches(0.8), Inches(1.9), Inches(11.733), Inches(5.1))
    caja28 = slide28.shapes.add_textbox(Inches(1.1), Inches(2.1), Inches(11.1), Inches(4.7))
    tf28 = caja28.text_frame
    tf28.word_wrap = True
    tf28.paragraphs[0].text = "Ganador de la Etapa: Supervivencia Aditiva"
    tf28.paragraphs[0].font.bold = True
    tf28.paragraphs[0].font.size = Pt(18)
    tf28.paragraphs[0].font.color.rgb = COLOR_VERDE

    pts28 = [
        "Supervivencia Aditiva garantiza elitismo estricto: el mejor fitness de la corrida nunca decrece.",
        "Permite reutilizar el caché de fitness de los padres sobrevivientes, ahorrando evaluaciones costosas de renderizado.",
        "Exclusiva con K = N arriesga retroceder si los hijos mutan desfavorablemente.",
        "Se adopta Supervivencia Aditiva."
    ]
    agregar_bullet_points(tf28, pts28, 14)

    # -------------------------------------------------------------
    # 29. ETAPA 6: FIGURAS - BENCHMARK DE RENDER
    # -------------------------------------------------------------
    slide29 = crear_slide_base(prs, "Etapa 6: Geometría de Genes y Rendimiento", "Tiempo de renderizado según el tipo de figura")
    agregar_card(slide29, Inches(0.8), Inches(1.9), Inches(5.7), Inches(5.1))
    caja29_txt = slide29.shapes.add_textbox(Inches(1.0), Inches(2.1), Inches(5.3), Inches(4.7))
    tf29_txt = caja29_txt.text_frame
    tf29_txt.word_wrap = True
    tf29_txt.paragraphs[0].text = "Costo Computacional por Tipo"
    tf29_txt.paragraphs[0].font.bold = True
    tf29_txt.paragraphs[0].font.size = Pt(16)
    pts29 = [
        "Triángulos: Polígonos de 3 vértices, rasterización simple y eficiente.",
        "Cuadriláteros y Pentágonos: Mayor cantidad de aristas, incrementando ligeramente el tiempo de render por cuadro.",
        "Óvalos: Rasterización elíptica basada en radios y rotación.",
        "A 100 figuras, todas operan entre 1 ms y 3 ms por renderizado."
    ]
    agregar_bullet_points(tf29_txt, pts29, 13)

    agregar_card(slide29, Inches(6.8), Inches(1.9), Inches(5.7), Inches(5.1))
    caja29_lbl = slide29.shapes.add_textbox(Inches(7.0), Inches(2.0), Inches(5.3), Inches(0.4))
    caja29_lbl.text_frame.paragraphs[0].text = "Tiempo de Renderizado para 100 Figuras (ms)"
    caja29_lbl.text_frame.paragraphs[0].font.bold = True
    caja29_lbl.text_frame.paragraphs[0].font.size = Pt(13)
    agregar_imagen_segura(slide29, root / "results/analysis/op_05_figuras_render/tiempo_render.png", Inches(7.0), Inches(2.5), width=Inches(5.3))

    # -------------------------------------------------------------
    # 30. ETAPA 6: FIGURAS - EVOLUCIÓN
    # -------------------------------------------------------------
    slide30 = crear_slide_base(prs, "Etapa 6: Figuras en Ciclo Evolutivo", "Convergencia de fitness a 100 genes")
    agregar_card(slide30, Inches(0.8), Inches(1.9), Inches(5.7), Inches(5.1))
    caja30_1 = slide30.shapes.add_textbox(Inches(1.0), Inches(2.0), Inches(5.3), Inches(0.4))
    caja30_1.text_frame.paragraphs[0].text = "Fitness Máximo vs Generación"
    caja30_1.text_frame.paragraphs[0].font.bold = True
    caja30_1.text_frame.paragraphs[0].font.size = Pt(13)
    agregar_imagen_segura(slide30, root / "results/analysis/full_06_figuras/fitness_maximo.png", Inches(1.0), Inches(2.5), width=Inches(5.3))

    agregar_card(slide30, Inches(6.8), Inches(1.9), Inches(5.7), Inches(5.1))
    caja30_2 = slide30.shapes.add_textbox(Inches(7.0), Inches(2.0), Inches(5.3), Inches(0.4))
    caja30_2.text_frame.paragraphs[0].text = "Tiempo por Generación (s)"
    caja30_2.text_frame.paragraphs[0].font.bold = True
    caja30_2.text_frame.paragraphs[0].font.size = Pt(13)
    agregar_imagen_segura(slide30, root / "results/analysis/full_06_figuras/tiempo_generacion.png", Inches(7.0), Inches(2.5), width=Inches(5.3))

    # -------------------------------------------------------------
    # 31. ETAPA 6: VISUAL (GIFs)
    # -------------------------------------------------------------
    slide31 = crear_slide_base(prs, "Etapa 6: Reconstrucción Visual", "Comparación de los 4 tipos de figura geométrica")
    var_fig = [
        ("Triángulo", "results/analysis/full_06_figuras/triangulo/img/bron500.gif"),
        ("Cuadrilátero", "results/analysis/full_06_figuras/cuadrilatero/img/bron500.gif"),
        ("Pentágono", "results/analysis/full_06_figuras/pentagono/img/bron500.gif"),
        ("Óvalo", "results/analysis/full_06_figuras/ovalo/img/bron500.gif")
    ]
    for i, (nombre, ruta) in enumerate(var_fig):
        x = Inches(0.8 + i * 2.95)
        agregar_card(slide31, x, Inches(1.9), Inches(2.8), Inches(5.1))
        caja_lbl = slide31.shapes.add_textbox(x + Inches(0.1), Inches(2.1), Inches(2.6), Inches(0.4))
        caja_lbl.text_frame.paragraphs[0].text = nombre
        caja_lbl.text_frame.paragraphs[0].font.bold = True
        caja_lbl.text_frame.paragraphs[0].font.size = Pt(12)
        agregar_imagen_segura(slide31, root / ruta, x + Inches(0.15), Inches(2.7), width=Inches(2.5))

    # -------------------------------------------------------------
    # 32. CONFIGURACIÓN ÓPTIMA
    # -------------------------------------------------------------
    slide32 = crear_slide_base(prs, "Configuración Óptima Recomendada", "Consolidación de las variantes ganadoras del hilo conductor")
    agregar_card(slide32, Inches(0.8), Inches(1.9), Inches(11.733), Inches(5.1))
    caja32 = slide32.shapes.add_textbox(Inches(1.1), Inches(2.1), Inches(11.1), Inches(4.7))
    tf32 = caja32.text_frame
    tf32.word_wrap = True
    tf32.paragraphs[0].text = "Receta Óptima del Motor Genético"
    tf32.paragraphs[0].font.bold = True
    tf32.paragraphs[0].font.size = Pt(18)
    tf32.paragraphs[0].font.color.rgb = COLOR_TITULO

    config_optima = [
        "Inicialización: Sesgo de Color (Bounding Box) — Acelera 5x la convergencia temprana.",
        "Selección: Torneo Determinístico (k=3 a 5) — Presión selectiva robusta e independiente de la escala de fitness.",
        "Cruza: Dos Puntos / Anular — Preserva esquemas y capas continuas de renderizado.",
        "Mutación: MultiGen (k=5 a 15) — Perturbaciones acotadas para ajuste fino sin destruir la estructura global.",
        "Supervivencia: Aditiva con brecha generacional moderada — Elitismo estricto y máximo ahorro de re-renderizado.",
        "Figura: Triángulos (para agilidad y flexibilidad poligonal) u Óvalos (para transiciones suaves)."
    ]
    agregar_bullet_points(tf32, config_optima, 13)

    # -------------------------------------------------------------
    # 33. CIERRE
    # -------------------------------------------------------------
    slide_fin = prs.slides.add_slide(prs.slide_layouts[6])
    fondo_fin = slide_fin.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), Inches(13.333), Inches(7.5))
    fondo_fin.fill.solid()
    fondo_fin.fill.fore_color.rgb = COLOR_TITULO
    fondo_fin.line.fill.background()

    caja_fin = slide_fin.shapes.add_textbox(Inches(1.5), Inches(2.5), Inches(10.333), Inches(2.5))
    tf_fin = caja_fin.text_frame
    tf_fin.word_wrap = True
    p_fin = tf_fin.paragraphs[0]
    p_fin.text = "¿Preguntas?"
    p_fin.font.name = "Arial"
    p_fin.font.size = Pt(44)
    p_fin.font.bold = True
    p_fin.font.color.rgb = RGBColor(255, 255, 255)
    p_fin.alignment = PP_ALIGN.CENTER

    p_fin2 = tf_fin.add_paragraph()
    p_fin2.text = "Sistemas de Inteligencia Artificial — ITBA"
    p_fin2.font.name = "Arial"
    p_fin2.font.size = Pt(18)
    p_fin2.font.color.rgb = RGBColor(190, 227, 248)
    p_fin2.alignment = PP_ALIGN.CENTER
    p_fin2.space_before = Pt(14)

    Path(ruta_salida).parent.mkdir(parents=True, exist_ok=True)
    prs.save(ruta_salida)
    print(f"Presentación generada exitosamente en: {ruta_salida}")


if __name__ == "__main__":
    crear_presentacion()
